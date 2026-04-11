"""
Dynamic Prompt Assembly Service — builds prompts at runtime from database content.
All criteria come from framework_nodes, scales, and form templates.
Reads evidence file content when possible.
"""
import os
import uuid
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.models.assessment_engine import (
    AssessmentInstance, AssessmentResponse, AssessmentEvidence,
    AssessmentScale, AssessmentScaleLevel, AssessmentFormTemplate, AiProduct,
)
from app.models.compliance_framework import ComplianceFramework
from app.models.framework_node import FrameworkNode


async def build_prompt(db: AsyncSession, instance_id: uuid.UUID, node_id: uuid.UUID, ai_product_id: uuid.UUID | None = None) -> dict:
    """Build system + user prompts dynamically from database content."""

    # Load instance with relationships
    inst = (await db.execute(
        select(AssessmentInstance)
        .options(selectinload(AssessmentInstance.assessed_entity), selectinload(AssessmentInstance.framework))
        .where(AssessmentInstance.id == instance_id)
    )).scalar_one_or_none()
    if not inst:
        raise ValueError("Assessment instance not found")

    fw = inst.framework

    # Load product if applicable
    product = None
    if ai_product_id:
        product = (await db.execute(select(AiProduct).where(AiProduct.id == ai_product_id))).scalar_one_or_none()

    # Load node
    node = (await db.execute(select(FrameworkNode).where(FrameworkNode.id == node_id))).scalar_one_or_none()
    if not node:
        raise ValueError("Node not found")

    # Load existing response
    q = select(AssessmentResponse).where(AssessmentResponse.instance_id == instance_id, AssessmentResponse.node_id == node_id)
    if ai_product_id:
        q = q.where(AssessmentResponse.ai_product_id == ai_product_id)
    else:
        q = q.where(AssessmentResponse.ai_product_id.is_(None))
    response = (await db.execute(q)).scalar_one_or_none()

    # Load evidence files
    evidence_files = []
    if response:
        evidence_files = (await db.execute(
            select(AssessmentEvidence).where(AssessmentEvidence.response_id == response.id)
        )).scalars().all()

    # Read evidence file contents
    evidence_contents = await _read_evidence_files(evidence_files)

    # Build prompts
    system_prompt = _build_system_prompt(fw)
    user_prompt = _build_user_prompt(
        fw=fw, node=node, response=response,
        evidence_files=evidence_files, evidence_contents=evidence_contents,
        entity_name=inst.assessed_entity.name if inst.assessed_entity else "Unknown",
        product=product,
    )

    return {"system_prompt": system_prompt, "user_prompt": user_prompt}


async def _read_evidence_files(evidence_files: list) -> dict[str, str]:
    """Read text content from evidence files where possible."""
    contents = {}
    for ef in evidence_files:
        if not ef.file_path or not os.path.exists(ef.file_path):
            continue
        ext = os.path.splitext(ef.file_name)[1].lower()
        try:
            if ext == ".pdf":
                contents[ef.file_name] = _read_pdf(ef.file_path)
            elif ext == ".docx":
                contents[ef.file_name] = _read_docx(ef.file_path)
            elif ext == ".xlsx":
                contents[ef.file_name] = _read_xlsx(ef.file_path)
            elif ext == ".pptx":
                contents[ef.file_name] = _read_pptx(ef.file_path)
            elif ext in (".txt", ".csv", ".json", ".md"):
                with open(ef.file_path, "r", encoding="utf-8", errors="ignore") as f:
                    contents[ef.file_name] = f.read()[:10000]
            # Images (.png, .jpg) are skipped — LLM can't process them as text
        except Exception:
            contents[ef.file_name] = f"[Error reading file: {ef.file_name}]"

    return contents


def _read_pdf(path: str) -> str:
    try:
        import PyPDF2
        text = ""
        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages[:20]:  # Max 20 pages
                text += page.extract_text() or ""
        return text[:15000]
    except ImportError:
        return "[PDF reading requires PyPDF2 — install it to enable document content extraction]"


def _read_docx(path: str) -> str:
    try:
        import docx
        doc = docx.Document(path)
        text = "\n".join(p.text for p in doc.paragraphs)
        return text[:15000]
    except ImportError:
        return "[DOCX reading requires python-docx — install it to enable document content extraction]"


def _read_xlsx(path: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)
        text = ""
        for sheet in wb.sheetnames[:5]:
            ws = wb[sheet]
            text += f"\n--- Sheet: {sheet} ---\n"
            for row in ws.iter_rows(max_row=100, values_only=True):
                text += " | ".join(str(c) if c is not None else "" for c in row) + "\n"
        return text[:15000]
    except ImportError:
        return "[XLSX reading requires openpyxl — install it to enable spreadsheet content extraction]"


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        text = ""
        for slide in prs.slides[:30]:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text_frame.text + "\n"
        return text[:15000]
    except ImportError:
        return "[PPTX reading requires python-pptx — install it to enable presentation content extraction]"


def _build_system_prompt(fw) -> str:
    return f"""You are an expert compliance assessor evaluating against the {fw.name} ({fw.abbreviation}) {fw.version or ''} framework.

Your task is to review the provided answer and supporting evidence documents for a specific control requirement, then provide:
1. A compliance determination
2. Detailed review feedback
3. Justification for your determination
4. Recommendations if not fully compliant

## Compliance Levels
- **Compliant**: The control requirement is fully met with sufficient evidence.
- **Semi-Compliant**: The control requirement is partially met — some evidence exists but there are gaps.
- **Non-Compliant**: The control requirement is not met — insufficient or no evidence provided.

## Rules
- Assess STRICTLY against the criteria described in the control requirement description and guidance.
- Read and analyze ALL provided evidence document content carefully.
- Base your assessment on BOTH the answer AND the evidence documents.
- If evidence is partial or ambiguous, determine Semi-Compliant.
- If no evidence is provided, determine Non-Compliant unless the answer clearly demonstrates compliance.
- Cite specific evidence documents and specific content from them in your reasoning.

## Document Validation Rules (CRITICAL)
For EACH evidence document, you MUST check:
1. **Document date**: Does the document contain a date? Is it recent and relevant? Extract the date if found.
2. **Approval status**: Is the document formally approved? Look for approval stamps, "Approved by:", sign-off blocks, "APPROVED" watermarks, approval statements, or references to formal approval.
3. **Signatures**: Does the document contain signatures, digital signature blocks, or stamp impressions?

A document that is NOT dated, NOT approved, or lacks signatures/stamps is considered WEAK evidence. For full compliance, evidence documents SHOULD be dated, approved, and signed. If key documents lack these, the compliance status should be reduced (e.g., from Compliant to Semi-Compliant).

## IMPORTANT — Output Style Rules
- Do NOT mention "the consultant", "the assessor", "the entity", or organization names in review_feedback, justification, or recommendations.
- Write in an impersonal, objective tone focused purely on the control requirement and the evidence.
- Use phrases like "The answer demonstrates...", "The evidence shows...", "The requirement is met because...", "To achieve compliance..."
- Do NOT start with "Based on the consultant's assessment..." or "The entity has..."

## Response Format
Respond ONLY with valid JSON in this exact format, no other text:
{{
  "compliance_status": "Compliant" | "Semi-Compliant" | "Non-Compliant",
  "review_feedback": "<objective review of the answer quality and evidence, 2-3 paragraphs, no entity/consultant references>",
  "justification": "<clear justification for the compliance determination, citing specific evidence, impersonal tone>",
  "recommendations": "<specific actionable recommendations to achieve full compliance, or 'N/A' if compliant, impersonal tone>",
  "confidence": <number 0-100>,
  "document_analysis": [
    {{
      "file_name": "<exact filename>",
      "document_date": "<extracted date in YYYY-MM-DD format or null if not found>",
      "is_approved": <true if approval evidence found, false otherwise>,
      "approved_by": "<who approved, or null>",
      "has_signature": <true if signature/stamp found, false otherwise>,
      "notes": "<brief note about document quality>"
    }}
  ]
}}"""


def _build_user_prompt(fw, node, response, evidence_files, evidence_contents, entity_name, product) -> str:
    prompt = "## Assessment Context\n\n"
    prompt += f"**Framework:** {fw.name} {fw.version or ''}\n"
    prompt += f"**Entity:** {entity_name}\n"
    if product:
        prompt += f"**AI Product:** {product.name}\n"
        if product.description:
            prompt += f"**Product Description:** {product.description}\n"
        if product.product_type:
            prompt += f"**Product Type:** {product.product_type}\n"

    prompt += f"\n## Control Requirement\n\n"
    prompt += f"**Reference:** {node.reference_code or ''}\n"
    prompt += f"**Name:** {node.name}\n"
    if node.description:
        prompt += f"**Description:** {node.description}\n"
    if node.guidance:
        prompt += f"\n**Assessment Guidance:** {node.guidance}\n"
    if hasattr(node, 'evidence_type') and node.evidence_type:
        prompt += f"\n**Expected Evidence Type:** {node.evidence_type}\n"
    if hasattr(node, 'acceptance_criteria') and node.acceptance_criteria:
        prompt += f"\n**Acceptance Criteria:**\n{node.acceptance_criteria}\n"
    if hasattr(node, 'acceptance_criteria_en') and node.acceptance_criteria_en:
        prompt += f"\n**Acceptance Criteria (EN):**\n{node.acceptance_criteria_en}\n"
    if hasattr(node, 'spec_references') and node.spec_references:
        prompt += f"\n**Specification References:** {node.spec_references}\n"
    if hasattr(node, 'maturity_level') and node.maturity_level is not None:
        prompt += f"\n**Maturity Level:** {node.maturity_level}\n"
    if hasattr(node, 'priority') and node.priority:
        prompt += f"**Priority:** {node.priority}\n"

    # Consultant's answer
    prompt += "\n## Consultant's Answer\n\n"
    if response and response.response_data:
        answer = response.response_data.get("answer", "")
        if answer:
            prompt += f"{answer}\n"
        else:
            prompt += "No answer provided by the consultant.\n"
    else:
        prompt += "No answer provided by the consultant.\n"

    # Evidence documents with content
    prompt += f"\n## Supporting Evidence ({len(evidence_files)} documents)\n\n"
    if not evidence_files:
        prompt += "No evidence documents have been uploaded.\n"
    else:
        for ef in evidence_files:
            prompt += f"### Document: {ef.file_name}\n"
            prompt += f"Type: {ef.file_type or 'unknown'} | Size: {ef.file_size or 0} bytes\n\n"
            content = evidence_contents.get(ef.file_name)
            if content:
                prompt += f"**Content:**\n{content}\n\n"
            else:
                ext = os.path.splitext(ef.file_name)[1].lower()
                if ext in (".png", ".jpg", ".jpeg"):
                    prompt += "[Image file — content cannot be extracted as text]\n\n"
                else:
                    prompt += "[File content could not be extracted]\n\n"

    prompt += "\n## Task\n\n"
    prompt += "Based on the control requirement, the consultant's answer, and the evidence documents above, "
    prompt += "determine the compliance status and provide your review. Respond with the JSON format specified in your instructions."

    return prompt
